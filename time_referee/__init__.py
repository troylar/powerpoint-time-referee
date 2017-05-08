from pptx import Presentation
from datetime import timedelta
import re
import time


class PowerPointManager():

    def __init__(self, **kwargs):
        self.path = kwargs.pop("Path")
        self.validate_file()
        self.load_slides()

    def validate_file(self):
        try:
            self.presentation = Presentation(self.path)

        except Exception as e:
            self.is_valid_file = False
            raise OSError("Error opening presentation {0}: {1}"
                          .format(self.path, e))
        self.is_valid_file = True

    def load_slides(self):
        self.slides = []
        for s in self.presentation.slides:
            self.slides.append(Slide(s))

    def time_to_str(self, time_in_sec):
        return time.strftime("%H:%M:%S", time.gmtime(time_in_sec))

    @property
    def total_time_in_seconds(self):
        return int(sum(s.time_in_seconds for s in self.slides))

    @property
    def total_time(self):
        return self.time_to_str(self.total_time_in_seconds)

    @property
    def total_slides(self):
        return len(self.slides)

    def apply_timings(self):
        t = 0
        for s in self.slides:
            ct = self.start_time.datetime + timedelta(seconds=t)
            cur_time = ct.strftime("%H:%M:%S")
            rem_time = self.time_to_str(self.total_time_in_seconds - t)
            s.update_timing(cur_time, rem_time)
            t = t + s.time_in_seconds

        self.actual_time = time.strftime("%H:%M:%S", time.gmtime(t))

    def save(self):
        print "Saving . . ."
        self.presentation.save(self.path)


class Slide():
    def __init__(self, slide):
        self.slide = slide
        self.time_in_seconds = 0
        self.note_text = ""
        if self.slide.has_notes_slide:
            self.note_text = \
                    re.sub(u"(\u2018|\u2019)", "'",
                           self.slide.notes_slide.notes_text_frame.text)

            self.note_text = re.sub(u"(\u201c|\u201d)", '"', self.note_text)
            self.note_text = re.sub(u"(\u2014)", '--', self.note_text)
            self.note_text = re.sub(u"(\u2026)", '. . .', self.note_text)

            m = re.search('^##\s?(\d+)', self.note_text, re.M)
            if m:
                self.time_in_seconds = int(m.group(1).strip())

    def update_timing(self, current_time, remaining_time):
        self.current_time = current_time
        self.remaining_time = remaining_time

        if self.slide.has_notes_slide:
            m = re.search('^#\sSTART\sOF\sSLIDE: (\d+\:\d+\:\d+)',
                          self.note_text, re.M)
            if m:
                if not m.group(1).strip() == current_time:
                    self.note_text = \
                        re.sub(r'# START OF SLIDE: \d+\:\d+\:\d+',
                               '# START OF SLIDE: {0}'.format(current_time),
                               self.note_text)
            else:
                self.note_text += \
                    '\n# START OF SLIDE: {0}'.format(current_time)

            m = re.search('^#\sREMAINING\sTIME: (\d+\:\d+\:\d+)',
                          self.note_text, re.M)
            if m:
                if not m.group(1).strip() == remaining_time:
                    self.note_text = \
                        re.sub(r'# REMAINING TIME: \d+\:\d+\:\d+',
                               '# REMAINING TIME: {0}'.format(remaining_time),
                               self.note_text)
            else:
                self.note_text += \
                    '\n# REMAINING TIME: {0}'.format(remaining_time)

        self.slide.notes_slide.notes_text_frame.text = self.note_text
